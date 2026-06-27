"""
Document Ingestor — PDF/DOCX/PPT to Segments
=============================================

Extracts structured segments from PDF, DOCX, and PPT files using PyMuPDF
(for PDF, already in our requirements) and python-docx / python-pptx
(optional — graceful degradation if not installed).

The segments are then fed into text_chunker.hierarchical_chunking() to
produce RAG-ready chunks. This brings chunkr-style document intelligence
to our existing knowledge/rag_engine.py WITHOUT requiring chunkr's full
Rust + Docker + GPU stack.

ADAPTER PATTERN (no conflict with existing code):
- This module is self-contained
- It DOES NOT modify knowledge/rag_engine.py
- It IS consumed by the new /api/v1/rag/ingest-pdf endpoint
- Optional dependencies degrade gracefully (PyMuPDF is the only required one)
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import List

from .text_chunker import (
    Chunk,
    ChunkingConfig,
    Segment,
    SegmentType,
    hierarchical_chunking,
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Optional Dependencies (Lazy imports — only loaded when first used)
# ---------------------------------------------------------------------------
# PyMuPDF, python-docx, and python-pptx are heavy imports (PyMuPDF alone
# takes ~5s and 120MB). We defer them to function-call time instead of
# module-import time. This makes `import document_ingestor` instant.

PDF_AVAILABLE = True  # Will be set to False on first use if import fails
DOCX_AVAILABLE = True
PPTX_AVAILABLE = True

_fitz = None  # Lazy-loaded PyMuPDF module
_docx = None  # Lazy-loaded python-docx module
_pptx = None  # Lazy-loaded python-pptx module


def _get_fitz():
    """Lazy-load PyMuPDF. Returns the fitz module or raises ImportError."""
    global _fitz, PDF_AVAILABLE
    if _fitz is not None:
        return _fitz
    try:
        import fitz  # type: ignore[import-untyped]

        _fitz = fitz
        return _fitz
    except ImportError as e:
        PDF_AVAILABLE = False
        raise ImportError(
            "PyMuPDF (fitz) is required for PDF ingestion. pip install pymupdf"
        ) from e


def _get_docx():
    """Lazy-load python-docx."""
    global _docx, DOCX_AVAILABLE
    if _docx is not None:
        return _docx
    try:
        import docx  # type: ignore[import-untyped]

        _docx = docx
        return _docx
    except ImportError as e:
        DOCX_AVAILABLE = False
        raise ImportError(
            "python-docx is required for DOCX ingestion. pip install python-docx"
        ) from e


def _get_pptx():
    """Lazy-load python-pptx."""
    global _pptx, PPTX_AVAILABLE
    if _pptx is not None:
        return _pptx
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        _pptx = Presentation
        return _pptx
    except ImportError as e:
        PPTX_AVAILABLE = False
        raise ImportError(
            "python-pptx is required for PPTX ingestion. pip install python-pptx"
        ) from e


def _detect_heading_level(text: str, font_size: float, default_size: float = 12.0) -> SegmentType:
    """Heuristic: detect if a text block is a Title, SectionHeader, or Text.

    Based on font size relative to the document's default body size.
    - > 1.8× default → Title
    - > 1.3× default → SectionHeader
    - else → Text
    """
    if font_size >= default_size * 1.8:
        return SegmentType.TITLE
    if font_size >= default_size * 1.3:
        return SegmentType.SECTION_HEADER
    return SegmentType.TEXT


def extract_segments_from_pdf(
    pdf_path: Path,
    ignore_headers_and_footers: bool = True,
) -> List[Segment]:
    """Extract segments from a PDF using PyMuPDF.

    Each text block on each page becomes a Segment. Font size is used to
    classify the segment as Title, SectionHeader, or Text. Images are
    detected as Picture segments (with bbox, no content text).

    Args:
        pdf_path: Path to the PDF file
        ignore_headers_and_footers: If True, skip the top 10% and bottom 10%
            of each page (typical header/footer zones)

    Returns:
        List of Segment objects in document order
    """
    if not PDF_AVAILABLE:
        raise ImportError("PyMuPDF (fitz) is required for PDF ingestion. pip install pymupdf")

    fitz = _get_fitz()  # Lazy-load PyMuPDF
    segments: List[Segment] = []
    doc = fitz.open(str(pdf_path))

    try:
        # First pass: compute the median font size to use as default_size
        font_sizes: List[float] = []
        for page in doc:
            blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
            for block in blocks:
                if block.get("type", 0) != 0:  # 0 = text block
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        if span.get("text", "").strip():
                            font_sizes.append(span.get("size", 12.0))
        default_size = sorted(font_sizes)[len(font_sizes) // 2] if font_sizes else 12.0

        # Second pass: extract segments
        for page_num, page in enumerate(doc, start=1):
            page_height = page.rect.height
            header_zone = page_height * 0.10 if ignore_headers_and_footers else 0
            footer_zone = page_height * 0.90 if ignore_headers_and_footers else page_height

            blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]

            for block in blocks:
                block_type = block.get("type", 0)
                bbox = block.get("bbox", (0, 0, 0, 0))

                # Image block → Picture segment
                if block_type == 1:  # 1 = image block
                    if ignore_headers_and_footers and (
                        bbox[1] < header_zone or bbox[3] > footer_zone
                    ):
                        continue
                    segments.append(
                        Segment(
                            content="[Image]",
                            segment_type=SegmentType.PICTURE,
                            page_number=page_num,
                            bbox=tuple(bbox),
                            metadata={"source_file": str(pdf_path)},
                        )
                    )
                    continue

                # Text block
                if block_type != 0:
                    continue

                # Skip header/footer zones
                if ignore_headers_and_footers and (bbox[1] < header_zone or bbox[3] > footer_zone):
                    # Classify as PageHeader or PageFooter for transparency
                    is_header = bbox[1] < header_zone
                    text = " ".join(
                        span.get("text", "")
                        for line in block.get("lines", [])
                        for span in line.get("spans", [])
                    ).strip()
                    if text:
                        segments.append(
                            Segment(
                                content=text,
                                segment_type=(
                                    SegmentType.PAGE_HEADER
                                    if is_header
                                    else SegmentType.PAGE_FOOTER
                                ),
                                page_number=page_num,
                                bbox=tuple(bbox),
                                metadata={"source_file": str(pdf_path)},
                            )
                        )
                    continue

                # Extract text + max font size from the block's spans
                block_text_parts: List[str] = []
                max_font_size = 0.0
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text", "")
                        if text.strip():
                            block_text_parts.append(text)
                            size = span.get("size", 12.0)
                            if size > max_font_size:
                                max_font_size = size

                block_text = " ".join(block_text_parts).strip()
                if not block_text:
                    continue

                # Detect tables (heuristic: blocks with many short lines + borders)
                # PyMuPDF doesn't expose table detection directly in get_text("dict"),
                # so we use a simple heuristic: if the block has > 5 lines and
                # the average line length is < 30 chars, classify as Table.
                line_count = len(block.get("lines", []))
                avg_line_len = len(block_text) / max(line_count, 1)
                if line_count > 5 and avg_line_len < 30:
                    segment_type = SegmentType.TABLE
                else:
                    segment_type = _detect_heading_level(block_text, max_font_size, default_size)

                segments.append(
                    Segment(
                        content=block_text,
                        segment_type=segment_type,
                        page_number=page_num,
                        bbox=tuple(bbox),
                        confidence=0.9,  # PyMuPDF extraction is reliable
                        metadata={
                            "source_file": str(pdf_path),
                            "font_size": max_font_size,
                        },
                    )
                )
    finally:
        doc.close()

    logger.info(
        "extract_segments_from_pdf: %s → %d segments (%d pages)",
        pdf_path.name,
        len(segments),
        len(doc) if PDF_AVAILABLE else 0,
    )
    return segments


def extract_segments_from_docx(docx_path: Path) -> List[Segment]:
    """Extract segments from a Word DOCX file using python-docx.

    Each paragraph becomes a Segment. Heading styles are mapped to
    Title / SectionHeader / Text.
    """
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx is required for DOCX ingestion. pip install python-docx")

    docx = _get_docx()  # Lazy-load python-docx
    segments: List[Segment] = []
    doc = docx.Document(str(docx_path))

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Map Word heading styles to our SegmentType
        style_name = (para.style.name or "").lower() if para.style else ""
        if "title" in style_name or "heading 1" in style_name:
            seg_type = SegmentType.TITLE
        elif "heading" in style_name:
            seg_type = SegmentType.SECTION_HEADER
        elif "caption" in style_name:
            seg_type = SegmentType.CAPTION
        else:
            seg_type = SegmentType.TEXT

        segments.append(
            Segment(
                content=text,
                segment_type=seg_type,
                metadata={"source_file": str(docx_path), "style": style_name},
            )
        )

    # Also extract tables
    for table in doc.tables:
        rows_text = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows_text.append(" | ".join(cells))
        if rows_text:
            segments.append(
                Segment(
                    content="\n".join(rows_text),
                    segment_type=SegmentType.TABLE,
                    metadata={"source_file": str(docx_path)},
                )
            )

    logger.info("extract_segments_from_docx: %s → %d segments", docx_path.name, len(segments))
    return segments


def extract_segments_from_pptx(pptx_path: Path) -> List[Segment]:
    """Extract segments from a PowerPoint PPTX file using python-pptx.

    Each slide's title becomes a SectionHeader; each text frame becomes Text.
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required for PPTX ingestion. pip install python-pptx")

    Presentation = _get_pptx()  # Lazy-load python-pptx
    segments: List[Segment] = []
    prs = Presentation(str(pptx_path))

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Title placeholder
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text_frame.text.strip()
            if title_text:
                segments.append(
                    Segment(
                        content=title_text,
                        segment_type=(
                            SegmentType.TITLE if slide_num == 1 else SegmentType.SECTION_HEADER
                        ),
                        page_number=slide_num,
                        metadata={"source_file": str(pptx_path)},
                    )
                )

        # Body text frames
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue  # Already added as title
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    segments.append(
                        Segment(
                            content=text,
                            segment_type=SegmentType.TEXT,
                            page_number=slide_num,
                            metadata={"source_file": str(pptx_path)},
                        )
                    )

        # Tables
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows_text = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_text.append(" | ".join(cells))
                if rows_text:
                    segments.append(
                        Segment(
                            content="\n".join(rows_text),
                            segment_type=SegmentType.TABLE,
                            page_number=slide_num,
                            metadata={"source_file": str(pptx_path)},
                        )
                    )

    logger.info("extract_segments_from_pptx: %s → %d segments", pptx_path.name, len(segments))
    return segments


def ingest_document(
    file_path: Path,
    target_length: int = 500,
    ignore_headers_and_footers: bool = True,
) -> List[Chunk]:
    """Extract segments from a document and chunk them hierarchically.

    This is the main entry point. It:
    1. Detects the file type (PDF/DOCX/PPTX)
    2. Extracts segments using the appropriate extractor
    3. Runs hierarchical_chunking to produce RAG-ready chunks

    Args:
        file_path: Path to the document (PDF/DOCX/PPTX)
        target_length: Target word count per chunk (default 500)
        ignore_headers_and_footers: Skip page headers/footers (default True)

    Returns:
        List of Chunk objects ready for embedding + storage

    Raises:
        ValueError: If file type is unsupported or required lib missing
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = file_path.suffix.lower()
    if ext == ".pdf":
        segments = extract_segments_from_pdf(file_path, ignore_headers_and_footers)
    elif ext == ".docx":
        segments = extract_segments_from_docx(file_path)
    elif ext == ".pptx":
        segments = extract_segments_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: .pdf, .docx, .pptx")

    if not segments:
        logger.warning("No segments extracted from %s", file_path)
        return []

    config = ChunkingConfig(
        target_length=target_length,
        ignore_headers_and_footers=ignore_headers_and_footers,
    )
    chunks = hierarchical_chunking(segments, config)

    logger.info(
        "ingest_document: %s → %d segments → %d chunks (target_length=%d words)",
        file_path.name,
        len(segments),
        len(chunks),
        target_length,
    )
    return chunks


def get_supported_formats() -> dict:
    """Return which document formats are currently supported.

    Useful for the API endpoint to tell the client what's available.
    """
    return {
        "pdf": {"supported": PDF_AVAILABLE, "library": "PyMuPDF (fitz)"},
        "docx": {"supported": DOCX_AVAILABLE, "library": "python-docx"},
        "pptx": {"supported": PPTX_AVAILABLE, "library": "python-pptx"},
    }


__all__ = [
    "extract_segments_from_pdf",
    "extract_segments_from_docx",
    "extract_segments_from_pptx",
    "ingest_document",
    "get_supported_formats",
]
